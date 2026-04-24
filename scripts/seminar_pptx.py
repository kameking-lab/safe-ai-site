"""
ANZEN AI セミナー PPTX 自動生成エンジン
====================================

YAML データを読み込んで、Midnight Executive 系デザインの 16:9 PPTX を生成する。

【設計方針】
- 単一ファイル・10レイアウト対応
- すべて Pt(...)/Inches(...) ベースで座標管理
- テキストはみ出しは「フィット幅」と「行送り」で制御
- 表紙とまとめのみネイビー背景、それ以外は白

【レイアウト一覧】
  cover          : 表紙（ネイビー）
  agenda         : 目次
  big_number     : 大数字 + 棒グラフ
  four_quadrant  : 4象限カード
  table          : 表
  flow           : 4ステップ横並び
  cards_2col     : 2カラム×4項目
  five_steps     : 5ステップ
  duties         : 4項目カード（管理者責務）
  takeaways      : まとめ（ネイビー）
"""

from __future__ import annotations

import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# =====================================================================
# デザインシステム
# =====================================================================
NAVY = RGBColor(0x1A, 0x27, 0x44)
NAVY_DEEP = RGBColor(0x12, 0x1B, 0x33)
GOLD = RGBColor(0xA0, 0x78, 0x3C)
GOLD_LIGHT = RGBColor(0xC9, 0xA8, 0x6E)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_JP = "Meiryo"
FONT_EN = "Meiryo"

# スライドサイズ 16:9（13.333" x 7.5"）
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# =====================================================================
# ユーティリティ
# =====================================================================
def set_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_no_line(shape) -> None:
    shape.line.fill.background()


def set_line(shape, rgb: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill: RGBColor | None = None,
             line_color: RGBColor | None = None, line_width_pt: float = 0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        set_fill(shape, fill)
    else:
        shape.fill.background()
    if line_color is not None:
        set_line(shape, line_color, line_width_pt)
    else:
        set_no_line(shape)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text: str, *,
             font_size: int = 14,
             font_color: RGBColor = SLATE_700,
             bold: bool = False,
             font_name: str = FONT_JP,
             align: int = PP_ALIGN.LEFT,
             anchor: int = MSO_ANCHOR.TOP,
             line_spacing: float | None = 1.2,
             tight: bool = False):
    """テキストボックスを配置。tight=True で内側余白を最小化。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if tight:
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
    else:
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color

    # eastAsia フォント明示（日本語の確実なレンダリング）
    rPr = run._r.get_or_add_rPr()
    for child in rPr.findall(qn("a:ea")):
        rPr.remove(child)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)

    return tb


def add_multiline(slide, x, y, w, h, lines: list[dict], *,
                  align: int = PP_ALIGN.LEFT,
                  anchor: int = MSO_ANCHOR.TOP,
                  tight: bool = False):
    """複数段落のテキストボックス。各 line は dict(text, size, color, bold, space_after, line_spacing)。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if tight:
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
    else:
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)

    first = True
    for ln in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if ln.get("line_spacing") is not None:
            p.line_spacing = ln["line_spacing"]
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        run = p.add_run()
        run.text = ln["text"]
        run.font.name = ln.get("font", FONT_JP)
        run.font.size = Pt(ln.get("size", 14))
        run.font.bold = ln.get("bold", False)
        run.font.color.rgb = ln.get("color", SLATE_700)

        rPr = run._r.get_or_add_rPr()
        for child in rPr.findall(qn("a:ea")):
            rPr.remove(child)
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", ln.get("font", FONT_JP))
    return tb


def add_bg(slide, color: RGBColor) -> None:
    """スライド全体の背景。"""
    add_rect(slide, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN), fill=color)


def add_title_block(slide, title_jp: str, title_en: str | None = None,
                    *, top_in: float = 0.55, left_in: float = 0.6):
    """白背景スライド用の標準タイトル。"""
    if title_en:
        add_text(
            slide,
            Inches(left_in), Inches(top_in),
            Inches(SLIDE_W_IN - left_in * 2), Inches(0.3),
            title_en, font_size=11, font_color=GOLD, bold=True, tight=True,
        )
        add_text(
            slide,
            Inches(left_in), Inches(top_in + 0.32),
            Inches(SLIDE_W_IN - left_in * 2), Inches(0.55),
            title_jp, font_size=26, font_color=NAVY, bold=True, tight=True,
        )
        # 下線（金）
        add_rect(
            slide,
            Inches(left_in), Inches(top_in + 0.95),
            Inches(0.6), Emu(28000), fill=GOLD,
        )
    else:
        add_text(
            slide,
            Inches(left_in), Inches(top_in),
            Inches(SLIDE_W_IN - left_in * 2), Inches(0.7),
            title_jp, font_size=26, font_color=NAVY, bold=True, tight=True,
        )
        add_rect(
            slide,
            Inches(left_in), Inches(top_in + 0.78),
            Inches(0.6), Emu(28000), fill=GOLD,
        )


def add_page_footer(slide, page_no: int, supervisor: str) -> None:
    """白背景の各スライド用ページ番号と監修者。"""
    add_text(
        slide,
        Inches(0.6), Inches(SLIDE_H_IN - 0.42),
        Inches(8.0), Inches(0.3),
        supervisor, font_size=8, font_color=SLATE_400, tight=True,
    )
    add_text(
        slide,
        Inches(SLIDE_W_IN - 1.6), Inches(SLIDE_H_IN - 0.42),
        Inches(1.0), Inches(0.3),
        f"{page_no:02d} / 10", font_size=9, font_color=SLATE_400,
        align=PP_ALIGN.RIGHT, bold=True, tight=True,
    )


# =====================================================================
# 各レイアウト
# =====================================================================
def render_cover(slide, data: dict, meta: dict) -> None:
    add_bg(slide, NAVY)

    # 装飾：左上の細い金ライン
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.6), Emu(28000), fill=GOLD)

    # バッジ
    add_text(
        slide,
        Inches(0.6), Inches(0.7),
        Inches(4.0), Inches(0.35),
        data.get("badge", ""), font_size=11, font_color=GOLD_LIGHT,
        bold=True, tight=True,
    )

    # 英タイトル
    add_text(
        slide,
        Inches(0.6), Inches(1.65),
        Inches(11.5), Inches(0.5),
        data.get("english_title", ""),
        font_size=14, font_color=GOLD, bold=True, tight=True,
    )

    # 日本語タイトル
    add_text(
        slide,
        Inches(0.6), Inches(2.15),
        Inches(11.5), Inches(1.4),
        data.get("title", ""),
        font_size=54, font_color=WHITE, bold=True, tight=True,
        line_spacing=1.05,
    )

    # サブタイトル
    add_text(
        slide,
        Inches(0.6), Inches(3.7),
        Inches(11.5), Inches(0.5),
        data.get("subtitle", ""),
        font_size=20, font_color=SLATE_300, tight=True,
    )

    # 区切り線
    add_rect(slide, Inches(0.6), Inches(5.3), Inches(2.0), Emu(15000), fill=GOLD)

    # 法的根拠 + 時間
    add_text(
        slide,
        Inches(0.6), Inches(5.45),
        Inches(11.5), Inches(0.4),
        data.get("legal", ""),
        font_size=13, font_color=SLATE_300, tight=True,
    )

    # 監修者（最下部）
    add_text(
        slide,
        Inches(0.6), Inches(SLIDE_H_IN - 0.65),
        Inches(11.5), Inches(0.35),
        data.get("supervisor", meta.get("supervisor", "")),
        font_size=10, font_color=SLATE_400, tight=True,
    )

    # 右下の装飾：ANZEN AI ロゴ風テキスト
    add_text(
        slide,
        Inches(SLIDE_W_IN - 2.4), Inches(0.7),
        Inches(1.8), Inches(0.35),
        "ANZEN AI",
        font_size=12, font_color=GOLD, bold=True,
        align=PP_ALIGN.RIGHT, tight=True,
    )


def render_agenda(slide, data: dict, meta: dict, page_no: int) -> None:
    add_title_block(slide, data["title"], data.get("title_en"))

    items = data["items"]
    # 2列 x 3行
    col_w = 5.85
    row_h = 1.5
    gap_x = 0.25
    gap_y = 0.18
    start_x = 0.6
    start_y = 1.95

    for i, item in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)

        # カード枠
        card = add_rect(
            slide,
            Inches(x), Inches(y),
            Inches(col_w), Inches(row_h),
            fill=WHITE, line_color=SLATE_300, line_width_pt=0.75,
        )

        # 番号（金）
        add_text(
            slide,
            Inches(x + 0.25), Inches(y + 0.18),
            Inches(0.7), Inches(0.5),
            item["num"], font_size=22, font_color=GOLD, bold=True, tight=True,
        )

        # 縦の細線
        add_rect(
            slide,
            Inches(x + 1.05), Inches(y + 0.25),
            Emu(10000), Inches(row_h - 0.5),
            fill=SLATE_300,
        )

        # 見出し
        add_text(
            slide,
            Inches(x + 1.2), Inches(y + 0.22),
            Inches(col_w - 1.4), Inches(0.5),
            item["head"], font_size=15, font_color=NAVY, bold=True, tight=True,
        )
        # 本文
        add_text(
            slide,
            Inches(x + 1.2), Inches(y + 0.78),
            Inches(col_w - 1.4), Inches(row_h - 0.85),
            item["body"], font_size=11, font_color=SLATE_700,
            line_spacing=1.35, tight=True,
        )

    add_page_footer(slide, page_no, meta.get("supervisor", ""))


def render_big_number(slide, data: dict, meta: dict, page_no: int) -> None:
    add_title_block(slide, data["title"], data.get("title_en"))

    # 左：大数字パネル
    panel_x = 0.6
    panel_y = 1.95
    panel_w = 5.5
    panel_h = 4.6
    add_rect(
        slide,
        Inches(panel_x), Inches(panel_y),
        Inches(panel_w), Inches(panel_h),
        fill=NAVY,
    )
    # 装飾の金ライン
    add_rect(
        slide,
        Inches(panel_x + 0.4), Inches(panel_y + 0.45),
        Inches(0.6), Emu(28000), fill=GOLD,
    )

    # 大数字
    add_multiline(
        slide,
        Inches(panel_x + 0.4), Inches(panel_y + 0.85),
        Inches(panel_w - 0.8), Inches(2.6),
        [
            {"text": str(data.get("big_value", "")),
             "size": 150, "color": WHITE, "bold": True, "line_spacing": 1.0},
        ],
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, tight=True,
    )
    # 単位
    add_text(
        slide,
        Inches(panel_x + 3.4), Inches(panel_y + 1.55),
        Inches(1.5), Inches(1.5),
        data.get("big_unit", ""),
        font_size=72, font_color=GOLD, bold=True, tight=True,
    )

    # キャプション
    add_text(
        slide,
        Inches(panel_x + 0.4), Inches(panel_y + 3.55),
        Inches(panel_w - 0.8), Inches(0.7),
        data.get("big_caption", ""),
        font_size=15, font_color=SLATE_300, line_spacing=1.35, tight=True,
    )
    # 出典
    add_text(
        slide,
        Inches(panel_x + 0.4), Inches(panel_y + panel_h - 0.5),
        Inches(panel_w - 0.8), Inches(0.35),
        f"出典：{data.get('source', '')}",
        font_size=9, font_color=SLATE_400, tight=True,
    )

    # 右：棒グラフ
    chart_x = 6.5
    chart_y = 1.95
    chart_w = 6.2
    chart_h = 4.6

    add_text(
        slide,
        Inches(chart_x), Inches(chart_y),
        Inches(chart_w), Inches(0.4),
        data.get("chart_title", ""), font_size=12,
        font_color=NAVY, bold=True, tight=True,
    )

    chart_data = data.get("chart_data", [])
    if chart_data:
        bar_top = chart_y + 0.55
        bar_bottom = chart_y + chart_h - 0.7
        bar_area_h = bar_bottom - bar_top
        n = len(chart_data)
        bar_h = min(0.5, (bar_area_h - (n - 1) * 0.18) / n)
        gap = (bar_area_h - bar_h * n) / max(n - 1, 1)

        max_val = max(item["value"] for item in chart_data)
        label_w = 1.95
        bar_x = chart_x + label_w + 0.05
        bar_max_w = chart_w - label_w - 0.15

        for i, item in enumerate(chart_data):
            y_i = bar_top + i * (bar_h + gap)
            # ラベル
            add_text(
                slide,
                Inches(chart_x), Inches(y_i + bar_h / 2 - 0.12),
                Inches(label_w - 0.05), Inches(0.3),
                item["label"], font_size=10, font_color=SLATE_700,
                align=PP_ALIGN.RIGHT, tight=True,
            )
            # 背景バー（薄)
            add_rect(
                slide,
                Inches(bar_x), Inches(y_i),
                Inches(bar_max_w), Inches(bar_h),
                fill=SLATE_100,
            )
            # 値バー
            ratio = item["value"] / max_val
            bar_w_use = bar_max_w * ratio
            color = NAVY if i == 0 else GOLD if i == 1 else SLATE_400
            add_rect(
                slide,
                Inches(bar_x), Inches(y_i),
                Inches(bar_w_use), Inches(bar_h),
                fill=color,
            )
            # 値ラベル
            add_text(
                slide,
                Inches(bar_x + bar_w_use + 0.05),
                Inches(y_i + bar_h / 2 - 0.12),
                Inches(1.0), Inches(0.3),
                f"{item['value']:.2f}",
                font_size=10, font_color=SLATE_700, bold=True, tight=True,
            )

    # 注記
    add_text(
        slide,
        Inches(chart_x), Inches(chart_y + chart_h - 0.4),
        Inches(chart_w), Inches(0.35),
        data.get("chart_note", ""),
        font_size=8, font_color=SLATE_500, tight=True,
    )

    add_page_footer(slide, page_no, meta.get("supervisor", ""))


def render_four_quadrant(slide, data: dict, meta: dict, page_no: int) -> None:
    add_title_block(slide, data["title"], data.get("title_en"))

    if data.get("intro"):
        add_text(
            slide,
            Inches(0.6), Inches(1.95),
            Inches(SLIDE_W_IN - 1.2), Inches(0.45),
            data["intro"], font_size=12, font_color=SLATE_700, tight=True,
        )

    # 2x2 グリッド
    quadrants = data["quadrants"]
    col_w = 5.85
    row_h = 2.05
    gap = 0.25
    start_x = 0.6
    start_y = 2.5

    for i, q in enumerate(quadrants):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_w + gap)
        y = start_y + row * (row_h + gap)

        # カード本体
        add_rect(
            slide,
            Inches(x), Inches(y),
            Inches(col_w), Inches(row_h),
            fill=SLATE_50, line_color=SLATE_300, line_width_pt=0.75,
        )
        # 左側の太いネイビーバー
        add_rect(
            slide,
            Inches(x), Inches(y),
            Inches(0.12), Inches(row_h),
            fill=NAVY,
        )

        # 番号
        add_text(
            slide,
            Inches(x + 0.32), Inches(y + 0.2),
            Inches(0.8), Inches(0.45),
            q["num"], font_size=20, font_color=GOLD, bold=True, tight=True,
        )

        # 見出し
        add_text(
            slide,
            Inches(x + 1.1), Inches(y + 0.22),
            Inches(col_w - 1.3), Inches(0.5),
            q["head"], font_size=17, font_color=NAVY, bold=True, tight=True,
        )

        # 本文
        add_text(
            slide,
            Inches(x + 0.32), Inches(y + 0.85),
            Inches(col_w - 0.55), Inches(row_h - 0.95),
            q["body"], font_size=10.5, font_color=SLATE_700,
            line_spacing=1.45, tight=True,
        )

    add_page_footer(slide, page_no, meta.get("supervisor", ""))


def render_table(slide, data: dict, meta: dict, page_no: int) -> None:
    add_title_block(slide, data["title"], data.get("title_en"))

    headers = data["headers"]
    rows = data["rows"]
    n_cols = len(headers)
    n_rows = len(rows) + 1  # ヘッダ含む

    # サイズ
    table_x = 0.6
    table_y = 2.0
    table_w = SLIDE_W_IN - 1.2
    table_h = 4.2

    # 列幅比率
    col_ratios = [0.18, 0.24, 0.28, 0.30]
    if len(col_ratios) != n_cols:
        col_ratios = [1.0 / n_cols] * n_cols
    col_widths = [table_w * r for r in col_ratios]

    header_h = 0.55
    row_h = (table_h - header_h) / (n_rows - 1)

    # ヘッダ
    cur_x = table_x
    for i, head in enumerate(headers):
        add_rect(
            slide,
            Inches(cur_x), Inches(table_y),
            Inches(col_widths[i]), Inches(header_h),
            fill=NAVY,
        )
        add_text(
            slide,
            Inches(cur_x + 0.1), Inches(table_y),
            Inches(col_widths[i] - 0.2), Inches(header_h),
            head, font_size=12, font_color=WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE, tight=True,
        )
        cur_x += col_widths[i]

    # 行
    for r_idx, row_data in enumerate(rows):
        cur_x = table_x
        y = table_y + header_h + r_idx * row_h
        bg = WHITE if r_idx % 2 == 0 else SLATE_50
        for c_idx, cell in enumerate(row_data):
            add_rect(
                slide,
                Inches(cur_x), Inches(y),
                Inches(col_widths[c_idx]), Inches(row_h),
                fill=bg, line_color=SLATE_300, line_width_pt=0.5,
            )
            color = NAVY if c_idx == 0 else SLATE_700
            bold = c_idx == 0
            add_text(
                slide,
                Inches(cur_x + 0.12), Inches(y),
                Inches(col_widths[c_idx] - 0.24), Inches(row_h),
                str(cell), font_size=11, font_color=color, bold=bold,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3, tight=True,
            )
            cur_x += col_widths[c_idx]

    # フットノート
    if data.get("footnote"):
        add_text(
            slide,
            Inches(table_x), Inches(table_y + table_h + 0.1),
            Inches(table_w), Inches(0.35),
            data["footnote"],
            font_size=9, font_color=SLATE_500, tight=True,
        )

    add_page_footer(slide, page_no, meta.get("supervisor", ""))


def render_flow(slide, data: dict, meta: dict, page_no: int) -> None:
    add_title_block(slide, data["title"], data.get("title_en"))

    if data.get("intro"):
        add_text(
            slide,
            Inches(0.6), Inches(1.95),
            Inches(SLIDE_W_IN - 1.2), Inches(0.4),
            data["intro"], font_size=12, font_color=SLATE_700, tight=True,
        )

    steps = data["steps"]
    n = len(steps)

    total_w = SLIDE_W_IN - 1.2
    gap = 0.2
    arrow_w = 0.45
    card_w = (total_w - (n - 1) * (gap + arrow_w)) / n
    card_h = 3.5
    start_x = 0.6
    start_y = 2.7

    for i, step in enumerate(steps):
        x = start_x + i * (card_w + gap + arrow_w)

        # カード
        add_rect(
            slide,
            Inches(x), Inches(start_y),
            Inches(card_w), Inches(card_h),
            fill=WHITE, line_color=SLATE_300, line_width_pt=0.75,
        )
        # 上部のネイビー帯
        add_rect(
            slide,
            Inches(x), Inches(start_y),
            Inches(card_w), Inches(0.7),
            fill=NAVY,
        )
        # STEP ラベル
        add_text(
            slide,
            Inches(x), Inches(start_y),
            Inches(card_w), Inches(0.7),
            step["num"], font_size=14, font_color=GOLD, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tight=True,
        )

        # 番号丸
        circle_size = 0.9
        circle_x = x + (card_w - circle_size) / 2
        circle_y = start_y + 0.95
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(circle_x), Inches(circle_y),
            Inches(circle_size), Inches(circle_size),
        )
        set_fill(circle, GOLD)
        set_no_line(circle)
        circle_tf = circle.text_frame
        circle_tf.margin_left = Pt(0)
        circle_tf.margin_right = Pt(0)
        circle_tf.margin_top = Pt(0)
        circle_tf.margin_bottom = Pt(0)
        circle_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = circle_tf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run()
        crun.text = str(i + 1)
        crun.font.name = FONT_EN
        crun.font.size = Pt(34)
        crun.font.bold = True
        crun.font.color.rgb = WHITE

        # 見出し
        add_text(
            slide,
            Inches(x + 0.15), Inches(start_y + 2.05),
            Inches(card_w - 0.3), Inches(0.55),
            step["head"], font_size=14, font_color=NAVY, bold=True,
            align=PP_ALIGN.CENTER, line_spacing=1.2, tight=True,
        )
        # 本文
        add_text(
            slide,
            Inches(x + 0.2), Inches(start_y + 2.65),
            Inches(card_w - 0.4), Inches(0.75),
            step["body"], font_size=10, font_color=SLATE_700,
            align=PP_ALIGN.CENTER, line_spacing=1.4, tight=True,
        )

        # 矢印
        if i < n - 1:
            ax = x + card_w + gap / 2
            ay = start_y + card_h / 2 - 0.2
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(ax), Inches(ay),
                Inches(arrow_w), Inches(0.4),
            )
            set_fill(arrow, SLATE_300)
            set_no_line(arrow)

    add_page_footer(slide, page_no, meta.get("supervisor", ""))


def render_cards_2col(slide, data: dict, meta: dict, page_no: int) -> None:
    add_title_block(slide, data["title"], data.get("title_en"))

    col_w = 5.95
    col_h = 4.55
    gap = 0.25
    start_x = 0.6
    start_y = 1.95

    for col_idx, side in enumerate(["left", "right"]):
        x = start_x + col_idx * (col_w + gap)
        head = data[f"{side}_head"]
        subhead = data.get(f"{side}_subhead", "")
        items = data[f"{side}_items"]

        # カラム枠
        add_rect(
            slide,
            Inches(x), Inches(start_y),
            Inches(col_w), Inches(col_h),
            fill=SLATE_50, line_color=SLATE_300, line_width_pt=0.75,
        )

        # ヘッダ帯（ネイビー）
        add_rect(
            slide,
            Inches(x), Inches(start_y),
            Inches(col_w), Inches(0.7),
            fill=NAVY,
        )
        # 英文小タイトル
        add_text(
            slide,
            Inches(x + 0.3), Inches(start_y + 0.05),
            Inches(col_w - 0.6), Inches(0.25),
            subhead, font_size=9, font_color=GOLD_LIGHT, bold=True,
            tight=True,
        )
        # 見出し
        add_text(
            slide,
            Inches(x + 0.3), Inches(start_y + 0.28),
            Inches(col_w - 0.6), Inches(0.4),
            head, font_size=15, font_color=WHITE, bold=True, tight=True,
        )

        # 4 アイテム
        item_top = start_y + 0.85
        item_h = (col_h - 0.95) / 4
        for j, it in enumerate(items):
            iy = item_top + j * item_h

            # 番号四角
            add_rect(
                slide,
                Inches(x + 0.25), Inches(iy + 0.1),
                Inches(0.4), Inches(0.4), fill=GOLD,
            )
            add_text(
                slide,
                Inches(x + 0.25), Inches(iy + 0.1),
                Inches(0.4), Inches(0.4),
                str(j + 1), font_size=14, font_color=WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tight=True,
            )
            # 見出し
            add_text(
                slide,
                Inches(x + 0.8), Inches(iy + 0.05),
                Inches(col_w - 1.0), Inches(0.4),
                it["head"], font_size=12, font_color=NAVY, bold=True,
                line_spacing=1.2, tight=True,
            )
            # 本文
            add_text(
                slide,
                Inches(x + 0.8), Inches(iy + 0.5),
                Inches(col_w - 1.0), Inches(item_h - 0.5),
                it["body"], font_size=10, font_color=SLATE_700,
                line_spacing=1.4, tight=True,
            )

    add_page_footer(slide, page_no, meta.get("supervisor", ""))


def render_five_steps(slide, data: dict, meta: dict, page_no: int) -> None:
    add_title_block(slide, data["title"], data.get("title_en"))

    if data.get("intro"):
        add_text(
            slide,
            Inches(0.6), Inches(1.95),
            Inches(SLIDE_W_IN - 1.2), Inches(0.4),
            data["intro"], font_size=12, font_color=SLATE_700, tight=True,
        )

    steps = data["steps"]
    n = len(steps)
    total_w = SLIDE_W_IN - 1.2
    gap = 0.18
    card_w = (total_w - (n - 1) * gap) / n
    card_h = 3.7
    start_x = 0.6
    start_y = 2.6

    for i, step in enumerate(steps):
        x = start_x + i * (card_w + gap)

        # カード
        add_rect(
            slide,
            Inches(x), Inches(start_y),
            Inches(card_w), Inches(card_h),
            fill=WHITE, line_color=SLATE_300, line_width_pt=0.75,
        )
        # 上部の金帯
        add_rect(
            slide,
            Inches(x), Inches(start_y),
            Inches(card_w), Inches(0.16), fill=GOLD,
        )

        # 番号（大）
        add_text(
            slide,
            Inches(x), Inches(start_y + 0.4),
            Inches(card_w), Inches(1.4),
            step["num"], font_size=64, font_color=NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tight=True,
            line_spacing=1.0,
        )

        # 見出し
        add_text(
            slide,
            Inches(x + 0.15), Inches(start_y + 1.85),
            Inches(card_w - 0.3), Inches(0.7),
            step["head"], font_size=14, font_color=NAVY, bold=True,
            align=PP_ALIGN.CENTER, line_spacing=1.25, tight=True,
        )
        # 本文
        add_text(
            slide,
            Inches(x + 0.18), Inches(start_y + 2.6),
            Inches(card_w - 0.36), Inches(card_h - 2.65),
            step["body"], font_size=10, font_color=SLATE_700,
            align=PP_ALIGN.CENTER, line_spacing=1.45, tight=True,
        )

    add_page_footer(slide, page_no, meta.get("supervisor", ""))


def render_duties(slide, data: dict, meta: dict, page_no: int) -> None:
    """4項目グリッド（管理者の責務）。four_quadrant の応用形だが
    背景は白＋ゴールドの番号アクセントで「責務」感を出す。"""
    add_title_block(slide, data["title"], data.get("title_en"))

    if data.get("intro"):
        add_text(
            slide,
            Inches(0.6), Inches(1.95),
            Inches(SLIDE_W_IN - 1.2), Inches(0.4),
            data["intro"], font_size=12, font_color=SLATE_700, tight=True,
        )

    items = data["items"]
    col_w = 5.85
    row_h = 2.0
    gap = 0.25
    start_x = 0.6
    start_y = 2.5

    for i, it in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_w + gap)
        y = start_y + row * (row_h + gap)

        # カード
        add_rect(
            slide,
            Inches(x), Inches(y),
            Inches(col_w), Inches(row_h),
            fill=WHITE, line_color=SLATE_300, line_width_pt=0.75,
        )
        # 番号バッジ（金）
        badge_w = 0.85
        add_rect(
            slide,
            Inches(x + 0.3), Inches(y + 0.3),
            Inches(badge_w), Inches(0.55), fill=GOLD,
        )
        add_text(
            slide,
            Inches(x + 0.3), Inches(y + 0.3),
            Inches(badge_w), Inches(0.55),
            it["num"], font_size=18, font_color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tight=True,
        )

        # 見出し
        add_text(
            slide,
            Inches(x + 1.3), Inches(y + 0.32),
            Inches(col_w - 1.5), Inches(0.5),
            it["head"], font_size=16, font_color=NAVY, bold=True, tight=True,
        )
        # 本文
        add_text(
            slide,
            Inches(x + 0.3), Inches(y + 0.95),
            Inches(col_w - 0.5), Inches(row_h - 1.05),
            it["body"], font_size=10.5, font_color=SLATE_700,
            line_spacing=1.45, tight=True,
        )

    add_page_footer(slide, page_no, meta.get("supervisor", ""))


def render_takeaways(slide, data: dict, meta: dict, page_no: int) -> None:
    add_bg(slide, NAVY)

    # タイトル
    add_text(
        slide,
        Inches(0.6), Inches(0.6),
        Inches(SLIDE_W_IN - 1.2), Inches(0.3),
        data.get("title_en", ""), font_size=11, font_color=GOLD,
        bold=True, tight=True,
    )
    add_text(
        slide,
        Inches(0.6), Inches(0.92),
        Inches(SLIDE_W_IN - 1.2), Inches(0.6),
        data["title"], font_size=28, font_color=WHITE, bold=True, tight=True,
    )
    # 金ライン
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(0.6), Emu(28000), fill=GOLD)

    # 4 takeaways
    items = data["items"]
    col_w = 5.95
    row_h = 1.95
    gap = 0.22
    start_x = 0.6
    start_y = 1.85

    for i, it in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_w + gap)
        y = start_y + row * (row_h + gap)

        # カード（白）
        add_rect(
            slide,
            Inches(x), Inches(y),
            Inches(col_w), Inches(row_h),
            fill=WHITE,
        )
        # 左の太い金バー
        add_rect(
            slide,
            Inches(x), Inches(y),
            Inches(0.12), Inches(row_h), fill=GOLD,
        )
        # 番号
        add_text(
            slide,
            Inches(x + 0.32), Inches(y + 0.18),
            Inches(0.7), Inches(0.45),
            it["num"], font_size=20, font_color=GOLD, bold=True, tight=True,
        )
        # 見出し
        add_text(
            slide,
            Inches(x + 1.05), Inches(y + 0.2),
            Inches(col_w - 1.25), Inches(0.55),
            it["head"], font_size=14, font_color=NAVY, bold=True,
            line_spacing=1.25, tight=True,
        )
        # 本文
        add_text(
            slide,
            Inches(x + 0.32), Inches(y + 0.85),
            Inches(col_w - 0.5), Inches(row_h - 0.95),
            it["body"], font_size=11, font_color=SLATE_700,
            line_spacing=1.45, tight=True,
        )

    # フッター部
    footer_y = 6.05
    add_rect(slide, Inches(0.6), Inches(footer_y), Inches(SLIDE_W_IN - 1.2), Emu(15000), fill=GOLD)

    add_text(
        slide,
        Inches(0.6), Inches(footer_y + 0.15),
        Inches(2.5), Inches(0.35),
        data.get("footer_label", "お問い合わせ"),
        font_size=11, font_color=GOLD, bold=True, tight=True,
    )
    add_text(
        slide,
        Inches(0.6), Inches(footer_y + 0.5),
        Inches(7.0), Inches(0.4),
        data.get("footer_value", ""),
        font_size=15, font_color=WHITE, bold=True, tight=True,
    )
    add_text(
        slide,
        Inches(0.6), Inches(SLIDE_H_IN - 0.55),
        Inches(SLIDE_W_IN - 1.2), Inches(0.35),
        data.get("footer_supervisor", meta.get("supervisor", "")),
        font_size=9, font_color=SLATE_400, tight=True,
    )


# =====================================================================
# ディスパッチャ
# =====================================================================
LAYOUTS = {
    "cover": render_cover,
    "agenda": render_agenda,
    "big_number": render_big_number,
    "four_quadrant": render_four_quadrant,
    "table": render_table,
    "flow": render_flow,
    "cards_2col": render_cards_2col,
    "five_steps": render_five_steps,
    "duties": render_duties,
    "takeaways": render_takeaways,
}


def build_pptx(yaml_path: Path, template_path: Path | None, output_path: Path) -> None:
    with yaml_path.open("r", encoding="utf-8") as f:
        spec = yaml.safe_load(f)

    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
        # テンプレートに含まれるサンプルスライドは削除
        sldIdLst = prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            sldIdLst.remove(sldId)
    else:
        prs = Presentation()

    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank_layout = prs.slide_layouts[6]  # 完全な空白レイアウト
    meta = spec.get("meta", {})

    for idx, slide_data in enumerate(spec["slides"], start=1):
        layout_key = slide_data["layout"]
        renderer = LAYOUTS.get(layout_key)
        if renderer is None:
            raise ValueError(f"Unknown layout: {layout_key}")
        slide = prs.slides.add_slide(blank_layout)

        if layout_key == "cover":
            renderer(slide, slide_data, meta)
        else:
            renderer(slide, slide_data, meta, idx)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[OK] Generated {output_path} ({len(spec['slides'])} slides)")


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python seminar_pptx.py <yaml_path> [template_path] [output_path]",
              file=sys.stderr)
        sys.exit(1)
    yaml_path = Path(sys.argv[1])
    template_path = Path(sys.argv[2]) if len(sys.argv) > 2 and sys.argv[2] else None
    output_path = Path(sys.argv[3]) if len(sys.argv) > 3 else Path(
        "dist/seminars") / yaml_path.with_suffix(".pptx").name
    build_pptx(yaml_path, template_path, output_path)


if __name__ == "__main__":
    main()
